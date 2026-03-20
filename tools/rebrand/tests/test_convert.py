#!/usr/bin/env python3
"""Tests for the rebrand conversion engine.

Creates minimal test .pptx files in memory and verifies
that font/colour/footer transformations apply correctly.
"""

import json
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# Add parent dir to path
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))

from convert import convert_deck, hex_to_rgb
from extract_rules import extract_text_runs, extract_colours, rgb_to_hex


def make_test_pptx(text="Hello World", font_name="Arial", font_size=24,
                   font_color="#003366", bold=False) -> str:
    """Create a minimal .pptx with one slide and one text box."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(font_color)

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name


def test_font_mapping():
    """Font family and size should change per mapping."""
    source = make_test_pptx(font_name="Arial", font_size=24, bold=True)

    mapping = {
        "fonts": {
            "Arial|24.0|bold": {
                "family": "Avenir Next",
                "size_pt": 20.0,
                "bold": True,
                "italic": False,
            }
        },
        "colours": {},
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as out:
        report = convert_deck(source, mapping, out.name)

    assert report.success, f"Conversion failed: {report.errors}"
    assert report.fonts_changed > 0, "No fonts were changed"

    # Verify output
    result = Presentation(out.name)
    for slide in result.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        assert run.font.name == "Avenir Next", f"Font not changed: {run.font.name}"
                        assert run.font.size == Pt(20), f"Size not changed: {run.font.size}"

    print("✓ test_font_mapping passed")


def test_colour_mapping():
    """Text colours should change per mapping."""
    source = make_test_pptx(font_color="#003366")

    mapping = {
        "fonts": {},
        "colours": {
            "#003366": "#1A5276",
        },
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as out:
        report = convert_deck(source, mapping, out.name)

    assert report.success, f"Conversion failed: {report.errors}"
    assert report.colours_changed > 0, "No colours were changed"

    # Verify output
    result = Presentation(out.name)
    for slide in result.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        hex_val = f"#{run.font.color.rgb}"
                        assert hex_val.upper() == "#1A5276", f"Colour not changed: {hex_val}"

    print("✓ test_colour_mapping passed")


def test_footer_replacement():
    """Footer text should be replaced."""
    source = make_test_pptx(text="© 2024 OldBrand Ltd")

    mapping = {
        "fonts": {},
        "colours": {},
        "footer": {
            "find": "© 2024 OldBrand Ltd",
            "replace": "© 2026 NewBrand Ltd",
        },
    }

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as out:
        report = convert_deck(source, mapping, out.name)

    assert report.success
    assert report.footers_updated > 0, "No footers were updated"

    # Verify output
    result = Presentation(out.name)
    for slide in result.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    if "Brand" in full_text:
                        assert "NewBrand" in full_text, f"Footer not replaced: {full_text}"

    print("✓ test_footer_replacement passed")


def test_no_mapping_no_changes():
    """Empty mapping should not change anything."""
    source = make_test_pptx()

    mapping = {"fonts": {}, "colours": {}}

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as out:
        report = convert_deck(source, mapping, out.name)

    assert report.success
    assert report.fonts_changed == 0
    assert report.colours_changed == 0

    print("✓ test_no_mapping_no_changes passed")


def test_extract_text_runs():
    """Text run extraction should capture font details."""
    source = make_test_pptx(text="Test text", font_name="Helvetica", font_size=18)
    prs = Presentation(source)
    runs = extract_text_runs(prs)

    assert len(runs) > 0, "No runs extracted"
    assert any(r["text"] == "Test text" for r in runs), "Text not found in runs"
    assert any(r["font_family"] == "Helvetica" for r in runs), "Font family not extracted"

    print("✓ test_extract_text_runs passed")


def test_extract_colours():
    """Colour extraction should find text colours."""
    source = make_test_pptx(font_color="#FF6600")
    prs = Presentation(source)
    colours = extract_colours(prs)

    assert "#FF6600" in colours, f"Colour not found. Found: {colours}"

    print("✓ test_extract_colours passed")


if __name__ == "__main__":
    test_font_mapping()
    test_colour_mapping()
    test_footer_replacement()
    test_no_mapping_no_changes()
    test_extract_text_runs()
    test_extract_colours()
    print("\n✓ All tests passed")

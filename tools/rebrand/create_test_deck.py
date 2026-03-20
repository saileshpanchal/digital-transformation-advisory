#!/usr/bin/env python3
"""Create a test input deck in the old brand style for conversion testing.

This simulates one of the client's 50 real decks — a different presentation
using the old brand (Meridian Consulting) visual identity.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path


def add_textbox(slide, left, top, width, height, text, font_name, font_size,
                font_color, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_bullets(slide, left, top, width, height, items, font_name, font_size, font_color):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Old brand colours
    NAVY = RGBColor(0x00, 0x33, 0x66)
    GOLD = RGBColor(0xCC, 0x99, 0x00)
    GREY = RGBColor(0x66, 0x66, 0x66)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    CREAM = RGBColor(0xF5, 0xF0, 0xE8)
    PURPLE = RGBColor(0x66, 0x66, 0x99)

    FOOTER = "© 2024 Meridian Consulting Group. All rights reserved. Confidential."

    # ── Slide 1: Title ───────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.6),
                "Meridian Consulting", "Georgia", Pt(22), WHITE)

    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(1.2),
                "Cloud Migration Assessment", "Georgia", Pt(44), WHITE, bold=True)

    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(8), Inches(0.8),
                "Current State Analysis & Migration Strategy", "Cambria", Pt(22), GOLD)

    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(8), Inches(0.5),
                "Prepared for Barclays Corporate Banking  |  February 2025", "Cambria", Pt(14), GREY, italic=True)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
                FOOTER, "Cambria", Pt(8), GREY)

    # ── Slide 2: Executive Summary ──────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), CREAM)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                "Executive Summary", "Georgia", Pt(32), WHITE, bold=True)

    add_bullets(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5), [
        "The organisation currently operates 147 on-premises applications across three data centres.",
        "Our assessment identifies 82 applications suitable for cloud migration within the next 18 months.",
        "We recommend a phased approach beginning with non-critical workloads to establish patterns and confidence.",
        "Total estimated cost reduction of 35% on infrastructure spend is achievable within 24 months.",
        "Critical success factors include executive sponsorship, cloud skills development, and network modernisation.",
    ], "Cambria", Pt(14), DARK)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                FOOTER, "Cambria", Pt(8), GREY)

    # ── Slide 3: Current Challenges ─────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), CREAM)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                "Current Challenges", "Georgia", Pt(32), WHITE, bold=True)

    # Three challenge cards
    challenges = [
        ("Legacy Infrastructure", "Ageing hardware with end-of-life dates approaching. Maintenance costs are escalating at 12% year-on-year with diminishing vendor support."),
        ("Talent Constraints", "The organisation is experiencing significant difficulty in recruiting and retaining skilled infrastructure engineers for on-premises environments."),
        ("Speed to Market", "New product features require infrastructure provisioning that takes 6-8 weeks. Competitors are deploying equivalent capabilities in days."),
    ]

    for i, (heading, body) in enumerate(challenges):
        left = Inches(0.8) + i * Inches(4.0)
        add_rect(slide, left, Inches(1.6), Inches(3.5), Inches(4.5), WHITE)
        # Card border
        shape = slide.shapes[-1]
        shape.line.color.rgb = PURPLE
        shape.line.width = Pt(1)

        add_textbox(slide, left + Inches(0.3), Inches(1.8), Inches(2.9), Inches(0.5),
                    heading, "Georgia", Pt(22), NAVY, bold=True)
        add_textbox(slide, left + Inches(0.3), Inches(2.5), Inches(2.9), Inches(3.0),
                    body, "Cambria", Pt(14), DARK)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                FOOTER, "Cambria", Pt(8), GREY)

    # ── Slide 4: Recommendation ─────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                "Our Recommendation", "Georgia", Pt(32), WHITE, bold=True)

    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.4),
                "A three-phase migration programme delivering sustainable cost reduction and operational agility",
                "Cambria", Pt(14), GOLD, italic=True)

    phases = [
        ("Phase 1: Foundation", "Q2 2025", "Establish cloud landing zone, governance framework, and security baseline. Migrate 15 non-critical workloads as pathfinder applications."),
        ("Phase 2: Accelerate", "Q4 2025", "Scale migration to 40 additional applications. Implement automated deployment pipelines and monitoring. Begin skills transformation programme."),
        ("Phase 3: Optimise", "Q2 2026", "Complete remaining migrations. Implement FinOps practices. Decommission legacy data centre capacity and realise cost savings."),
    ]

    for i, (name, timeline, desc) in enumerate(phases):
        top = Inches(2.2) + i * Inches(1.6)
        add_textbox(slide, Inches(0.8), top, Inches(3), Inches(0.4),
                    name, "Georgia", Pt(22), GOLD, bold=True)
        add_textbox(slide, Inches(4), top, Inches(2), Inches(0.4),
                    timeline, "Cambria", Pt(14), GREY, bold=True)
        add_textbox(slide, Inches(0.8), top + Inches(0.5), Inches(10), Inches(0.8),
                    desc, "Cambria", Pt(14), WHITE)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                FOOTER, "Cambria", Pt(8), GREY)

    # ── Slide 5: Investment ─────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), CREAM)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                "Investment Summary", "Georgia", Pt(32), WHITE, bold=True)

    # Metrics
    metrics = [
        ("£2.4M", "Total programme investment"),
        ("£840K", "Annual savings from Year 1"),
        ("2.8 years", "Payback period"),
        ("187%", "3-year ROI"),
    ]

    for i, (value, label) in enumerate(metrics):
        left = Inches(0.8) + i * Inches(3.1)
        add_textbox(slide, left, Inches(1.8), Inches(2.8), Inches(0.8),
                    value, "Georgia", Pt(48), NAVY, bold=True)
        add_textbox(slide, left, Inches(2.7), Inches(2.8), Inches(0.4),
                    label, "Cambria", Pt(14), DARK)

    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(10), Inches(2.0),
                "The investment is structured across three financial years with the majority of expenditure "
                "concentrated in Year 1. We anticipate the programme will be self-funding by the end of Year 2, "
                "with cumulative savings exceeding investment by Q3 of Year 3. These projections are based on "
                "conservative assumptions and exclude potential revenue uplift from faster time-to-market.",
                "Cambria", Pt(14), DARK)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                FOOTER, "Cambria", Pt(8), GREY)

    # ── Slide 6: Next Steps ─────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                "Next Steps", "Georgia", Pt(32), WHITE, bold=True)

    add_bullets(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(3.0), [
        "Board approval of migration strategy and investment case",
        "Cloud landing zone design workshop: 3-day session with architecture team",
        "Vendor selection for managed cloud services partner",
        "Pathfinder application selection and migration sprint planning",
        "Establishment of Cloud Centre of Excellence with seconded client staff",
    ], "Cambria", Pt(14), WHITE)

    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(3), Inches(0.4),
                "Contact", "Georgia", Pt(22), GOLD, bold=True)
    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(5), Inches(0.3),
                "James Robertson, Engagement Director", "Cambria", Pt(14), WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(6.1), Inches(5), Inches(0.3),
                "j.robertson@example.com", "Cambria", Pt(14), GOLD)

    add_textbox(slide, Inches(9), Inches(6.8), Inches(4), Inches(0.4),
                "Meridian Consulting", "Georgia", Pt(22), WHITE, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(9), Inches(7.1), Inches(4), Inches(0.3),
                "Trusted advisors since 1998", "Cambria", Pt(11), GREY, italic=True, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                FOOTER, "Cambria", Pt(8), GREY)

    # Save
    out = Path(__file__).parent / "input" / "cloud-migration-barclays.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Created test deck: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""Create designer exemplar decks for the /rebrand skill.

Generates source-exemplar.pptx (old brand) and target-exemplar.pptx (new brand)
with identical text content but different visual treatment.

Old brand: "Meridian Consulting" — corporate blue, serif fonts, formal tone
New brand: "Apex Partners"       — modern teal/coral, sans-serif, conversational tone

Both decks have 6 slides:
  1. Title slide
  2. About / company overview
  3. Our approach (3-column)
  4. Case study metrics
  5. Team slide
  6. Contact / closing slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path


# ── Brand Definitions ──────────────────────────────────────────

OLD_BRAND = {
    "name": "Meridian Consulting",
    "tagline": "Trusted advisors since 1998",
    "heading_font": "Georgia",
    "body_font": "Cambria",
    "heading_size": Pt(32),
    "subheading_size": Pt(22),
    "body_size": Pt(14),
    "small_size": Pt(11),
    "primary": RGBColor(0x00, 0x33, 0x66),      # Dark navy
    "secondary": RGBColor(0x66, 0x66, 0x99),     # Muted purple-blue
    "accent": RGBColor(0xCC, 0x99, 0x00),         # Gold
    "bg_dark": RGBColor(0x00, 0x33, 0x66),        # Navy bg
    "bg_light": RGBColor(0xF5, 0xF0, 0xE8),       # Warm cream
    "text_dark": RGBColor(0x1A, 0x1A, 0x2E),      # Near black
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),      # White
    "text_muted": RGBColor(0x66, 0x66, 0x66),     # Grey
    "footer": "© 2024 Meridian Consulting Group. All rights reserved. Confidential.",
    "heading_bold": True,
    "body_bold": False,
}

NEW_BRAND = {
    "name": "Apex Partners",
    "tagline": "Strategy that moves",
    "heading_font": "Calibri",
    "body_font": "Calibri",
    "heading_size": Pt(36),
    "subheading_size": Pt(20),
    "body_size": Pt(13),
    "small_size": Pt(10),
    "primary": RGBColor(0x00, 0x7C, 0x7A),       # Deep teal
    "secondary": RGBColor(0xFF, 0x6B, 0x6B),      # Coral
    "accent": RGBColor(0x2E, 0xCC, 0x71),          # Green
    "bg_dark": RGBColor(0x1A, 0x1A, 0x2E),        # Dark charcoal
    "bg_light": RGBColor(0xFA, 0xFA, 0xFA),        # Clean white-grey
    "text_dark": RGBColor(0x2D, 0x2D, 0x2D),      # Soft black
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),      # White
    "text_muted": RGBColor(0x88, 0x88, 0x88),     # Light grey
    "footer": "Apex Partners Ltd. Private & Confidential.",
    "heading_bold": True,
    "body_bold": False,
}


# ── Slide Content (identical for both) ─────────────────────────

SLIDES = [
    {
        "type": "title",
        "title": "Digital Transformation Programme",
        "subtitle": "Strategic Roadmap & Delivery Plan",
        "detail": "Prepared for GlobalTech Industries  |  Q1 2025",
    },
    {
        "type": "overview",
        "title": "Who We Are",
        "bullets": [
            "We are a specialist technology strategy consultancy working with FTSE 250 and Fortune 500 organisations.",
            "Our team brings deep expertise in cloud migration, platform modernisation, and AI-native architecture.",
            "We have delivered over 200 transformation programmes across financial services, healthcare, and manufacturing.",
            "Our methodology is evidence-based, outcome-driven, and designed for sustainable change.",
        ],
        # Tone variants — same meaning, different voice
        "bullets_new_tone": [
            "We're a technology strategy firm that works with the world's most ambitious organisations.",
            "Our people bring hands-on expertise in cloud, platforms, and AI — not just slide decks.",
            "200+ transformation programmes delivered. Financial services, healthcare, manufacturing.",
            "Evidence-based. Outcome-driven. Built to last.",
        ],
    },
    {
        "type": "three_column",
        "title": "Our Approach",
        "columns": [
            {
                "heading": "Discover",
                "body": "Comprehensive assessment of your current technology landscape, organisational capabilities, and strategic objectives to establish a clear baseline.",
            },
            {
                "heading": "Design",
                "body": "Architecture blueprints, migration strategies, and implementation roadmaps tailored to your specific constraints and aspirations.",
            },
            {
                "heading": "Deliver",
                "body": "Hands-on programme delivery with embedded teams, continuous measurement, and adaptive planning to ensure sustainable outcomes.",
            },
        ],
        "columns_new_tone": [
            {
                "heading": "Discover",
                "body": "We map your tech landscape, team capabilities, and goals — building a clear picture of where you stand and where you're headed.",
            },
            {
                "heading": "Design",
                "body": "Blueprints, migration plans, and roadmaps shaped around your real constraints — not generic frameworks.",
            },
            {
                "heading": "Deliver",
                "body": "We embed with your teams, measure what matters, and adapt as we go. Outcomes over outputs.",
            },
        ],
    },
    {
        "type": "metrics",
        "title": "Case Study: NatWest Group",
        "subtitle": "Cloud-first platform migration — 18-month programme",
        "metrics": [
            ("40%", "Reduction in infrastructure costs"),
            ("3x", "Faster time-to-market for new products"),
            ("99.95%", "Platform availability achieved"),
            ("£12M", "Annual savings realised in Year 1"),
        ],
        "quote": "The programme exceeded our expectations on every dimension. The team's deep expertise in platform engineering made the difference.",
        "quote_attribution": "— Chief Technology Officer, NatWest Group",
    },
    {
        "type": "team",
        "title": "Your Engagement Team",
        "members": [
            ("Sarah Chen", "Programme Director", "20 years in digital transformation. Former VP Engineering at Barclays."),
            ("Marcus Wright", "Lead Architect", "Cloud-native specialist. AWS/Azure certified. 15 years in platform engineering."),
            ("Priya Sharma", "Delivery Manager", "Agile practitioner. Delivered £50M+ programmes across financial services."),
        ],
    },
    {
        "type": "closing",
        "title": "Next Steps",
        "bullets": [
            "Discovery workshop: 2-day on-site session with your leadership team",
            "Current state assessment: 4-week deep-dive into technology and organisation",
            "Strategic roadmap: Delivered within 8 weeks of engagement start",
            "Board-ready presentation: Executive summary for governance sign-off",
        ],
        "contact_name": "Sarah Chen",
        "contact_email": "sarah.chen@example.com",
        "contact_phone": "+44 20 7946 0958",
    },
]


# ── Helper Functions ───────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text, font_name, font_size,
                font_color, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with a single formatted run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_name, font_size,
                    font_color, bold=False):
    """Add a text box with bullet-pointed paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.level = 0
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold

    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ── Slide Builders ─────────────────────────────────────────────

def build_title_slide(prs, content, brand):
    """Slide 1: Title slide with brand colours and company name."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background band
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), brand["bg_dark"])

    # Accent stripe
    add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), brand["accent"])

    # Company name
    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.6),
                brand["name"], brand["heading_font"], brand["subheading_size"],
                brand["text_light"], bold=False, alignment=PP_ALIGN.LEFT)

    # Title
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(1.2),
                content["title"], brand["heading_font"], Pt(44),
                brand["text_light"], bold=brand["heading_bold"], alignment=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(8), Inches(0.8),
                content["subtitle"], brand["body_font"], brand["subheading_size"],
                brand["accent"], bold=False, alignment=PP_ALIGN.LEFT)

    # Detail line
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(8), Inches(0.5),
                content["detail"], brand["body_font"], brand["body_size"],
                brand["text_muted"], italic=True, alignment=PP_ALIGN.LEFT)

    # Tagline
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(6), Inches(0.5),
                brand["tagline"], brand["body_font"], brand["small_size"],
                brand["text_muted"], italic=True, alignment=PP_ALIGN.LEFT)

    # Footer
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
                brand["footer"], brand["body_font"], Pt(8),
                brand["text_muted"], alignment=PP_ALIGN.LEFT)


def build_overview_slide(prs, content, brand, use_new_tone=False):
    """Slide 2: Company overview with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light background
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), brand["bg_light"])

    # Coloured header bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), brand["primary"])

    # Title on header bar
    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                content["title"], brand["heading_font"], brand["heading_size"],
                brand["text_light"], bold=brand["heading_bold"], alignment=PP_ALIGN.LEFT)

    # Bullets
    bullets = content.get("bullets_new_tone") if use_new_tone else content["bullets"]
    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(4.5),
                    bullets, brand["body_font"], brand["body_size"],
                    brand["text_dark"], bold=brand["body_bold"])

    # Company name bottom-right
    add_textbox(slide, Inches(9), Inches(7.0), Inches(4), Inches(0.3),
                brand["name"], brand["body_font"], Pt(8),
                brand["text_muted"], alignment=PP_ALIGN.RIGHT)

    # Footer
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                brand["footer"], brand["body_font"], Pt(8),
                brand["text_muted"], alignment=PP_ALIGN.LEFT)


def build_three_column_slide(prs, content, brand, use_new_tone=False):
    """Slide 3: Three-column approach slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), brand["bg_light"])

    # Header bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), brand["primary"])

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                content["title"], brand["heading_font"], brand["heading_size"],
                brand["text_light"], bold=brand["heading_bold"])

    # Three columns
    columns = content.get("columns_new_tone") if use_new_tone else content["columns"]
    col_width = Inches(3.5)
    col_gap = Inches(0.4)
    start_left = Inches(0.8)

    for i, col in enumerate(columns):
        left = start_left + i * (col_width + col_gap)

        # Column card background
        add_rect(slide, left, Inches(1.6), col_width, Inches(4.5),
                 brand["text_light"], brand["secondary"])

        # Column number circle (using accent colour)
        add_textbox(slide, left + Inches(0.3), Inches(1.8), Inches(0.6), Inches(0.5),
                    str(i + 1), brand["heading_font"], Pt(24),
                    brand["accent"], bold=True, alignment=PP_ALIGN.CENTER)

        # Column heading
        add_textbox(slide, left + Inches(0.3), Inches(2.4), col_width - Inches(0.6), Inches(0.5),
                    col["heading"], brand["heading_font"], brand["subheading_size"],
                    brand["primary"], bold=True)

        # Column body
        add_textbox(slide, left + Inches(0.3), Inches(3.1), col_width - Inches(0.6), Inches(2.5),
                    col["body"], brand["body_font"], brand["body_size"],
                    brand["text_dark"])

    # Footer
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                brand["footer"], brand["body_font"], Pt(8),
                brand["text_muted"])


def build_metrics_slide(prs, content, brand):
    """Slide 4: Case study with big metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), brand["bg_dark"])

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                content["title"], brand["heading_font"], brand["heading_size"],
                brand["text_light"], bold=brand["heading_bold"])

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
                content["subtitle"], brand["body_font"], brand["body_size"],
                brand["accent"], italic=True)

    # Metrics in a 2x2 grid
    positions = [
        (Inches(0.8), Inches(2.2)),
        (Inches(4.8), Inches(2.2)),
        (Inches(0.8), Inches(4.0)),
        (Inches(4.8), Inches(4.0)),
    ]

    for (left, top), (value, label) in zip(positions, content["metrics"]):
        # Big number
        add_textbox(slide, left, top, Inches(3.5), Inches(0.8),
                    value, brand["heading_font"], Pt(48),
                    brand["accent"], bold=True)
        # Label
        add_textbox(slide, left, top + Inches(0.8), Inches(3.5), Inches(0.4),
                    label, brand["body_font"], brand["body_size"],
                    brand["text_light"])

    # Quote
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.6),
                f'"{content["quote"]}"', brand["body_font"], brand["body_size"],
                brand["text_light"], italic=True)

    # Attribution
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(10), Inches(0.3),
                content["quote_attribution"], brand["body_font"], brand["small_size"],
                brand["text_muted"])

    # Footer
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                brand["footer"], brand["body_font"], Pt(8),
                brand["text_muted"])


def build_team_slide(prs, content, brand):
    """Slide 5: Team members."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), brand["bg_light"])

    # Header bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), brand["primary"])

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                content["title"], brand["heading_font"], brand["heading_size"],
                brand["text_light"], bold=brand["heading_bold"])

    # Team members in columns
    col_width = Inches(3.5)
    start_left = Inches(0.8)

    for i, (name, role, bio) in enumerate(content["members"]):
        left = start_left + i * (col_width + Inches(0.4))

        # Card
        add_rect(slide, left, Inches(1.6), col_width, Inches(4.0),
                 brand["text_light"], brand["secondary"])

        # Name
        add_textbox(slide, left + Inches(0.3), Inches(1.8), col_width - Inches(0.6), Inches(0.4),
                    name, brand["heading_font"], brand["subheading_size"],
                    brand["primary"], bold=True)

        # Role
        add_textbox(slide, left + Inches(0.3), Inches(2.3), col_width - Inches(0.6), Inches(0.3),
                    role, brand["body_font"], brand["body_size"],
                    brand["accent"], bold=True)

        # Bio
        add_textbox(slide, left + Inches(0.3), Inches(2.8), col_width - Inches(0.6), Inches(2.5),
                    bio, brand["body_font"], brand["body_size"],
                    brand["text_dark"])

    # Footer
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                brand["footer"], brand["body_font"], Pt(8),
                brand["text_muted"])


def build_closing_slide(prs, content, brand):
    """Slide 6: Next steps & contact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), brand["bg_dark"])

    # Accent stripe
    add_rect(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), brand["accent"])

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                content["title"], brand["heading_font"], brand["heading_size"],
                brand["text_light"], bold=brand["heading_bold"])

    # Bullets
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(3.5),
                    content["bullets"], brand["body_font"], brand["body_size"],
                    brand["text_light"])

    # Contact section
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(3), Inches(0.4),
                "Contact", brand["heading_font"], brand["subheading_size"],
                brand["accent"], bold=True)

    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(5), Inches(0.3),
                content["contact_name"], brand["body_font"], brand["body_size"],
                brand["text_light"], bold=True)

    add_textbox(slide, Inches(0.8), Inches(6.1), Inches(5), Inches(0.3),
                content["contact_email"], brand["body_font"], brand["body_size"],
                brand["accent"])

    add_textbox(slide, Inches(0.8), Inches(6.4), Inches(5), Inches(0.3),
                content["contact_phone"], brand["body_font"], brand["body_size"],
                brand["text_light"])

    # Company name
    add_textbox(slide, Inches(9), Inches(6.8), Inches(4), Inches(0.4),
                brand["name"], brand["heading_font"], brand["subheading_size"],
                brand["text_light"], bold=False, alignment=PP_ALIGN.RIGHT)

    # Tagline
    add_textbox(slide, Inches(9), Inches(7.1), Inches(4), Inches(0.3),
                brand["tagline"], brand["body_font"], brand["small_size"],
                brand["text_muted"], italic=True, alignment=PP_ALIGN.RIGHT)

    # Footer
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                brand["footer"], brand["body_font"], Pt(8),
                brand["text_muted"])


# ── Main ───────────────────────────────────────────────────────

def create_exemplar(brand, slides_content, output_path, use_new_tone=False):
    """Create a complete exemplar deck with the given brand."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    for content in slides_content:
        slide_type = content["type"]
        if slide_type == "title":
            build_title_slide(prs, content, brand)
        elif slide_type == "overview":
            build_overview_slide(prs, content, brand, use_new_tone)
        elif slide_type == "three_column":
            build_three_column_slide(prs, content, brand, use_new_tone)
        elif slide_type == "metrics":
            build_metrics_slide(prs, content, brand)
        elif slide_type == "team":
            build_team_slide(prs, content, brand)
        elif slide_type == "closing":
            build_closing_slide(prs, content, brand)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"  Created: {output_path} ({len(prs.slides)} slides)")


def main():
    assets_dir = Path(__file__).parent / "assets"

    print("Creating designer exemplar decks...\n")

    # Source exemplar — old brand, formal tone
    create_exemplar(
        OLD_BRAND, SLIDES,
        str(assets_dir / "source-exemplar.pptx"),
        use_new_tone=False,
    )

    # Target exemplar — new brand, conversational tone
    create_exemplar(
        NEW_BRAND, SLIDES,
        str(assets_dir / "target-exemplar.pptx"),
        use_new_tone=True,
    )

    print(f"\nDone. Exemplar pair ready in {assets_dir}/")
    print("\nBrand transformation summary:")
    print(f"  Old brand:  {OLD_BRAND['name']} (Georgia/Cambria, navy/gold, formal)")
    print(f"  New brand:  {NEW_BRAND['name']} (Calibri, teal/coral, conversational)")
    print(f"\nNext: Run '/rebrand extract' to generate brand-mapping.json")


if __name__ == "__main__":
    main()

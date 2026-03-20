#!/usr/bin/env python3
"""Extract brand mapping rules from a pair of designer exemplar PowerPoint decks.

Compares source (old brand) and target (new brand) slide-by-slide to extract:
- Font mappings (family, size, bold, italic, colour)
- Colour palette mappings (hex → hex)
- Layout structure mappings (slide dimensions, margins, placeholders)
- Logo/image replacements
- Footer/header text changes

Usage:
    python extract_rules.py \
        --source assets/source-exemplar.pptx \
        --target assets/target-exemplar.pptx \
        --output brand-mapping.json

The designer must create both decks with identical content on matching slides,
so the script can diff the visual treatment of the same text/structure.
"""

import argparse
import json
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def rgb_to_hex(rgb_color):
    """Convert an RGBColor to a hex string."""
    if rgb_color is None:
        return None
    return f"#{rgb_color}"


def pt_to_float(pt_value):
    """Convert Pt EMU value to float points."""
    if pt_value is None:
        return None
    return pt_value.pt


def extract_text_runs(presentation):
    """Extract all text runs with their formatting from a presentation.

    Returns a list of dicts, each describing a text run with its
    slide index, shape name, paragraph index, run index, text content,
    and full formatting (font family, size, bold, italic, colour, alignment).
    """
    runs = []
    for slide_idx, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                alignment = str(para.alignment) if para.alignment else None
                for run_idx, run in enumerate(para.runs):
                    font = run.font
                    runs.append({
                        "slide": slide_idx,
                        "shape": shape.name,
                        "para": para_idx,
                        "run": run_idx,
                        "text": run.text,
                        "font_family": font.name,
                        "font_size_pt": pt_to_float(font.size),
                        "bold": font.bold,
                        "italic": font.italic,
                        "underline": font.underline,
                        "color": rgb_to_hex(font.color.rgb) if font.color and font.color.rgb else None,
                        "alignment": alignment,
                    })
    return runs


def extract_colours(presentation):
    """Extract all unique colours used in the presentation.

    Scans text runs, shape fills, shape borders, and background fills.
    Returns a set of hex colour strings.
    """
    colours = set()

    for slide in presentation.slides:
        # Background
        try:
            bg = slide.background
            if bg.fill and bg.fill.fore_color and bg.fill.fore_color.rgb:
                colours.add(rgb_to_hex(bg.fill.fore_color.rgb))
        except (TypeError, AttributeError):
            pass  # No fill or theme-based fill

        for shape in slide.shapes:
            # Shape fill
            if hasattr(shape, "fill"):
                try:
                    fill = shape.fill
                    if fill.fore_color and fill.fore_color.rgb:
                        colours.add(rgb_to_hex(fill.fore_color.rgb))
                except Exception:
                    pass

            # Shape line/border
            if hasattr(shape, "line"):
                try:
                    line = shape.line
                    if line.color and line.color.rgb:
                        colours.add(rgb_to_hex(line.color.rgb))
                except Exception:
                    pass

            # Text colours
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                colours.add(rgb_to_hex(run.font.color.rgb))
                        except Exception:
                            pass

    colours.discard(None)
    return colours


def extract_slide_layouts(presentation):
    """Extract slide layout info: dimensions, master names, layout names."""
    width = presentation.slide_width
    height = presentation.slide_height
    layouts = []
    for slide_idx, slide in enumerate(presentation.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        master_name = slide.slide_layout.slide_master.name if slide.slide_layout and slide.slide_layout.slide_master else "unknown"
        layouts.append({
            "slide": slide_idx,
            "layout": layout_name,
            "master": master_name,
        })
    return {
        "width_emu": width,
        "height_emu": height,
        "width_inches": width / 914400,
        "height_inches": height / 914400,
        "slides": layouts,
    }


def extract_images(presentation):
    """Extract image metadata (name, dimensions, content type) per slide."""
    images = []
    for slide_idx, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                images.append({
                    "slide": slide_idx,
                    "name": shape.name,
                    "width_emu": shape.width,
                    "height_emu": shape.height,
                    "left_emu": shape.left,
                    "top_emu": shape.top,
                    "content_type": shape.image.content_type if hasattr(shape, "image") else None,
                })
    return images


def build_font_mapping(source_runs, target_runs):
    """Match text runs between source and target by text content.

    For each unique text string found in both, compare the font formatting
    and build a mapping of source formatting → target formatting.
    """
    # Index target runs by (slide, text) for matching
    target_by_key = {}
    for run in target_runs:
        key = (run["slide"], run["text"].strip())
        if key not in target_by_key and run["text"].strip():
            target_by_key[key] = run

    font_map = {}
    for src_run in source_runs:
        key = (src_run["slide"], src_run["text"].strip())
        if key not in target_by_key or not src_run["text"].strip():
            continue

        tgt_run = target_by_key[key]

        src_font_key = (
            src_run["font_family"],
            src_run["font_size_pt"],
            src_run["bold"],
        )
        tgt_font_spec = {
            "family": tgt_run["font_family"],
            "size_pt": tgt_run["font_size_pt"],
            "bold": tgt_run["bold"],
            "italic": tgt_run["italic"],
        }

        if src_font_key not in font_map:
            font_map[src_font_key] = tgt_font_spec

    # Convert to serializable format
    result = {}
    for (family, size, bold), target in font_map.items():
        key_str = f"{family or 'inherit'}|{size or 'inherit'}|{'bold' if bold else 'normal'}"
        result[key_str] = target

    return result


def build_colour_mapping(source_runs, target_runs):
    """Match colours between source and target by text content position.

    Returns a hex→hex mapping of old colours to new colours.
    """
    target_by_key = {}
    for run in target_runs:
        key = (run["slide"], run["text"].strip())
        if key not in target_by_key and run["text"].strip():
            target_by_key[key] = run

    colour_map = {}
    for src_run in source_runs:
        key = (src_run["slide"], src_run["text"].strip())
        if key not in target_by_key:
            continue

        src_colour = src_run["color"]
        tgt_colour = target_by_key[key]["color"]

        if src_colour and tgt_colour and src_colour != tgt_colour:
            colour_map[src_colour] = tgt_colour

    return colour_map


def build_text_changes(source_runs, target_runs):
    """Find text that changed between source and target (tone/language changes).

    These are cases where the text itself was rewritten, not just reformatted.
    """
    # Index by (slide, shape, para, run) for positional matching
    target_by_pos = {}
    for run in target_runs:
        key = (run["slide"], run["shape"], run["para"], run["run"])
        target_by_pos[key] = run

    changes = []
    for src_run in source_runs:
        key = (src_run["slide"], src_run["shape"], src_run["para"], src_run["run"])
        if key not in target_by_pos:
            continue

        tgt_run = target_by_pos[key]
        if src_run["text"].strip() != tgt_run["text"].strip():
            if src_run["text"].strip() and tgt_run["text"].strip():
                changes.append({
                    "slide": src_run["slide"],
                    "source_text": src_run["text"],
                    "target_text": tgt_run["text"],
                })

    return changes


def extract_brand_mapping(source_path: str, target_path: str) -> dict:
    """Main extraction: compare source and target exemplar decks.

    Returns a complete brand mapping dictionary.
    """
    source = Presentation(source_path)
    target = Presentation(target_path)

    source_runs = extract_text_runs(source)
    target_runs = extract_text_runs(target)

    source_layouts = extract_slide_layouts(source)
    target_layouts = extract_slide_layouts(target)

    source_images = extract_images(source)
    target_images = extract_images(target)

    source_colours = extract_colours(source)
    target_colours = extract_colours(target)

    font_mapping = build_font_mapping(source_runs, target_runs)
    colour_mapping = build_colour_mapping(source_runs, target_runs)
    text_changes = build_text_changes(source_runs, target_runs)

    mapping = {
        "_meta": {
            "source_file": str(source_path),
            "target_file": str(target_path),
            "source_slide_count": len(source.slides),
            "target_slide_count": len(target.slides),
            "extraction_version": "1.0.0",
        },
        "dimensions": {
            "source": {
                "width_inches": source_layouts["width_inches"],
                "height_inches": source_layouts["height_inches"],
            },
            "target": {
                "width_inches": target_layouts["width_inches"],
                "height_inches": target_layouts["height_inches"],
            },
        },
        "fonts": font_mapping,
        "colours": colour_mapping,
        "source_colours_all": sorted(list(source_colours)),
        "target_colours_all": sorted(list(target_colours)),
        "text_changes": text_changes,
        "source_layouts": source_layouts["slides"],
        "target_layouts": target_layouts["slides"],
        "source_images": source_images,
        "target_images": target_images,
        "tone_rules": [],  # Populated by Claude after reviewing text_changes
        "logo_replacement": {
            "enabled": False,
            "source_name_pattern": "Logo",
            "target_image_path": "assets/new-logo.png",
        },
        "footer": {
            "find": "",
            "replace": "",
        },
    }

    return mapping


def main():
    parser = argparse.ArgumentParser(
        description="Extract brand mapping rules from source and target exemplar PowerPoint decks."
    )
    parser.add_argument("--source", required=True, help="Path to source (old brand) exemplar .pptx")
    parser.add_argument("--target", required=True, help="Path to target (new brand) exemplar .pptx")
    parser.add_argument("--output", default="brand-mapping.json", help="Output path for brand mapping JSON")

    args = parser.parse_args()

    if not Path(args.source).exists():
        print(f"ERROR: Source file not found: {args.source}", file=sys.stderr)
        sys.exit(1)
    if not Path(args.target).exists():
        print(f"ERROR: Target file not found: {args.target}", file=sys.stderr)
        sys.exit(1)

    print(f"Extracting brand rules...")
    print(f"  Source: {args.source}")
    print(f"  Target: {args.target}")

    mapping = extract_brand_mapping(args.source, args.target)

    with open(args.output, "w") as f:
        json.dump(mapping, f, indent=2, default=str)

    # Summary
    print(f"\nExtracted brand mapping to {args.output}:")
    print(f"  Font rules:    {len(mapping['fonts'])}")
    print(f"  Colour rules:  {len(mapping['colours'])}")
    print(f"  Text changes:  {len(mapping['text_changes'])}")
    print(f"  Source colours: {len(mapping['source_colours_all'])}")
    print(f"  Target colours: {len(mapping['target_colours_all'])}")
    print(f"  Source images:  {len(mapping['source_images'])}")
    print(f"  Target images:  {len(mapping['target_images'])}")

    if mapping["text_changes"]:
        print(f"\n  ⚠ {len(mapping['text_changes'])} text changes detected.")
        print(f"    Review brand-mapping.json > text_changes, then populate tone_rules.")
        print(f"    Run: claude 'Review the text_changes in brand-mapping.json and generate tone_rules'")

    if not mapping["fonts"]:
        print(f"\n  ⚠ No font mappings detected. Ensure both decks have identical text content on matching slides.")

    print(f"\nNext steps:")
    print(f"  1. Review brand-mapping.json — verify font/colour mappings look correct")
    print(f"  2. Add tone_rules (or let Claude generate them from text_changes)")
    print(f"  3. Set logo_replacement.target_image_path if swapping logos")
    print(f"  4. Set footer.find / footer.replace for footer text")
    print(f"  5. Run: python convert.py --input input/old-deck.pptx --mapping brand-mapping.json")


if __name__ == "__main__":
    main()

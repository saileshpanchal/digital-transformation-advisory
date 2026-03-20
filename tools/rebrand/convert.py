#!/usr/bin/env python3
"""Convert a PowerPoint deck from one brand to another using a brand mapping.

Applies font, colour, layout, logo, and footer transformations deterministically.
Text tone adjustments are handled by Claude via the skill (not this script).

Usage:
    # Single file
    python convert.py \
        --input input/old-deck.pptx \
        --mapping brand-mapping.json \
        --output output/new-deck.pptx

    # Batch (all .pptx in input/)
    python convert.py \
        --batch input/ \
        --mapping brand-mapping.json \
        --output output/
"""

import argparse
import copy
import json
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor


@dataclass
class ConversionReport:
    """Track what was changed in a single deck conversion."""
    input_file: str = ""
    output_file: str = ""
    slides_processed: int = 0
    fonts_changed: int = 0
    colours_changed: int = 0
    logos_replaced: int = 0
    footers_updated: int = 0
    images_flagged: int = 0
    warnings: list = field(default_factory=list)
    errors: list = field(default_factory=list)

    @property
    def success(self) -> bool:
        return len(self.errors) == 0

    def summary(self) -> str:
        status = "✓" if self.success else "✗"
        lines = [
            f"{status} {self.input_file} → {self.output_file}",
            f"  Slides: {self.slides_processed}",
            f"  Fonts changed: {self.fonts_changed}",
            f"  Colours changed: {self.colours_changed}",
            f"  Logos replaced: {self.logos_replaced}",
            f"  Footers updated: {self.footers_updated}",
        ]
        if self.images_flagged:
            lines.append(f"  ⚠ Images flagged for review: {self.images_flagged}")
        for w in self.warnings:
            lines.append(f"  ⚠ {w}")
        for e in self.errors:
            lines.append(f"  ✗ {e}")
        return "\n".join(lines)


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' to an RGBColor."""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def parse_font_key(key: str) -> tuple:
    """Parse 'Arial|12.0|bold' → (family, size, is_bold)."""
    parts = key.split("|")
    family = parts[0] if parts[0] != "inherit" else None
    size = float(parts[1]) if len(parts) > 1 and parts[1] != "inherit" else None
    bold = parts[2] == "bold" if len(parts) > 2 else False
    return family, size, bold


def apply_font_mapping(run, font_map: dict, report: ConversionReport):
    """Apply font mapping to a single text run."""
    font = run.font
    current_key = f"{font.name or 'inherit'}|{font.size.pt if font.size else 'inherit'}|{'bold' if font.bold else 'normal'}"

    if current_key in font_map:
        target = font_map[current_key]
        if target.get("family") and target["family"] != font.name:
            font.name = target["family"]
            report.fonts_changed += 1
        if target.get("size_pt") and (not font.size or target["size_pt"] != font.size.pt):
            font.size = Pt(target["size_pt"])
            report.fonts_changed += 1
        if target.get("bold") is not None and target["bold"] != font.bold:
            font.bold = target["bold"]
        if target.get("italic") is not None and target["italic"] != font.italic:
            font.italic = target["italic"]
        return

    # Fuzzy match: try without bold/normal distinction
    for map_key, target in font_map.items():
        src_family, src_size, src_bold = parse_font_key(map_key)
        if src_family and font.name and src_family.lower() == font.name.lower():
            if target.get("family"):
                font.name = target["family"]
                report.fonts_changed += 1
            if target.get("size_pt") and src_size and font.size:
                # Scale proportionally if sizes differ
                scale = target["size_pt"] / src_size
                font.size = Pt(round(font.size.pt * scale, 1))
                report.fonts_changed += 1
            break


def apply_colour_to_run(run, colour_map: dict, report: ConversionReport):
    """Apply colour mapping to a text run's font colour."""
    try:
        if run.font.color and run.font.color.rgb:
            hex_current = f"#{run.font.color.rgb}".upper()
            if hex_current in colour_map:
                run.font.color.rgb = hex_to_rgb(colour_map[hex_current])
                report.colours_changed += 1
    except Exception:
        pass  # Some colour types (theme, inherited) can't be read directly


def apply_colour_to_fill(fill, colour_map: dict, report: ConversionReport):
    """Apply colour mapping to a shape fill."""
    try:
        if fill.fore_color and fill.fore_color.rgb:
            hex_current = f"#{fill.fore_color.rgb}".upper()
            if hex_current in colour_map:
                fill.fore_color.rgb = hex_to_rgb(colour_map[hex_current])
                report.colours_changed += 1
    except Exception:
        pass


def apply_colour_to_line(line, colour_map: dict, report: ConversionReport):
    """Apply colour mapping to a shape line/border."""
    try:
        if line.color and line.color.rgb:
            hex_current = f"#{line.color.rgb}".upper()
            if hex_current in colour_map:
                line.color.rgb = hex_to_rgb(colour_map[hex_current])
                report.colours_changed += 1
    except Exception:
        pass


def replace_logo(shape, mapping: dict, report: ConversionReport) -> bool:
    """Replace a logo image if it matches the source pattern."""
    logo_config = mapping.get("logo_replacement", {})
    if not logo_config.get("enabled"):
        return False

    pattern = logo_config.get("source_name_pattern", "")
    target_path = logo_config.get("target_image_path", "")

    if not pattern or not target_path:
        return False

    if pattern.lower() in shape.name.lower():
        target = Path(target_path)
        if target.exists():
            # python-pptx doesn't support in-place image replacement easily,
            # so we replace the image blob in the relationship
            try:
                image_part = shape.image
                with open(target, "rb") as f:
                    image_part._blob = f.read()
                report.logos_replaced += 1
                return True
            except Exception as e:
                report.warnings.append(f"Logo replacement failed on '{shape.name}': {e}")
        else:
            report.warnings.append(f"Logo target file not found: {target_path}")

    return False


def update_footer_text(shape, mapping: dict, report: ConversionReport):
    """Replace footer text if it matches the find pattern."""
    footer_config = mapping.get("footer", {})
    find_text = footer_config.get("find", "")
    replace_text = footer_config.get("replace", "")

    if not find_text or not replace_text:
        return

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if find_text in run.text:
                    run.text = run.text.replace(find_text, replace_text)
                    report.footers_updated += 1


def convert_deck(input_path: str, mapping: dict, output_path: str) -> ConversionReport:
    """Convert a single PowerPoint deck using the brand mapping.

    Applies all deterministic transformations:
    - Font family/size/weight mapping
    - Colour hex mapping (text, fills, borders)
    - Logo replacement
    - Footer text replacement

    Does NOT handle tone/text rewriting — that's Claude's job via the skill.
    """
    report = ConversionReport(input_file=input_path, output_file=output_path)

    try:
        prs = Presentation(input_path)
    except Exception as e:
        report.errors.append(f"Failed to open: {e}")
        return report

    font_map = mapping.get("fonts", {})
    # Normalise colour map keys to uppercase for case-insensitive matching
    colour_map = {k.upper(): v for k, v in mapping.get("colours", {}).items()}

    for slide_idx, slide in enumerate(prs.slides):
        report.slides_processed += 1

        for shape in slide.shapes:
            # Logo replacement (check before other processing)
            if shape.shape_type == 13:  # Picture
                if replace_logo(shape, mapping, report):
                    continue
                # Flag non-logo images for review
                report.images_flagged += 1
                report.warnings.append(
                    f"Slide {slide_idx + 1}: Image '{shape.name}' flagged for manual review"
                )

            # Text formatting
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        apply_font_mapping(run, font_map, report)
                        apply_colour_to_run(run, colour_map, report)

                # Footer replacement
                update_footer_text(shape, mapping, report)

            # Shape fill colours
            if hasattr(shape, "fill"):
                apply_colour_to_fill(shape.fill, colour_map, report)

            # Shape border colours
            if hasattr(shape, "line"):
                apply_colour_to_line(shape.line, colour_map, report)

    # Save
    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
    except Exception as e:
        report.errors.append(f"Failed to save: {e}")

    return report


def batch_convert(input_dir: str, mapping: dict, output_dir: str) -> list[ConversionReport]:
    """Convert all .pptx files in input_dir."""
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    pptx_files = sorted(input_path.glob("*.pptx"))
    if not pptx_files:
        print(f"No .pptx files found in {input_dir}")
        return []

    reports = []
    for pptx_file in pptx_files:
        if pptx_file.name.startswith("~$"):  # Skip temp files
            continue
        out_file = output_path / pptx_file.name
        print(f"  Converting: {pptx_file.name}...")
        report = convert_deck(str(pptx_file), mapping, str(out_file))
        reports.append(report)
        print(f"    {'✓' if report.success else '✗'} {report.fonts_changed} fonts, {report.colours_changed} colours")

    return reports


def print_batch_report(reports: list[ConversionReport]):
    """Print a summary of batch conversion results."""
    total = len(reports)
    success = sum(1 for r in reports if r.success)
    total_fonts = sum(r.fonts_changed for r in reports)
    total_colours = sum(r.colours_changed for r in reports)
    total_logos = sum(r.logos_replaced for r in reports)
    total_warnings = sum(len(r.warnings) for r in reports)

    print(f"\n{'=' * 60}")
    print(f"BATCH CONVERSION REPORT")
    print(f"{'=' * 60}")
    print(f"  Files:    {success}/{total} converted successfully")
    print(f"  Fonts:    {total_fonts} text runs updated")
    print(f"  Colours:  {total_colours} colour values updated")
    print(f"  Logos:    {total_logos} logos replaced")
    if total_warnings:
        print(f"  Warnings: {total_warnings} items flagged for review")
    print(f"{'=' * 60}")

    # Per-file detail
    for report in reports:
        if not report.success or report.warnings:
            print(f"\n{report.summary()}")

    if total_warnings:
        print(f"\n⚠ {total_warnings} items need manual review. Check flagged images and warnings above.")

    needs_tone = any(True for r in reports if r.slides_processed > 0)
    if needs_tone:
        print(f"\nNext: Run '/rebrand tone input/' to apply tone adjustments via Claude.")


def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint decks from one brand to another."
    )
    parser.add_argument("--input", help="Path to a single .pptx file to convert")
    parser.add_argument("--batch", help="Path to directory of .pptx files to convert")
    parser.add_argument("--mapping", required=True, help="Path to brand-mapping.json")
    parser.add_argument("--output", required=True, help="Output file path (single) or directory (batch)")

    args = parser.parse_args()

    if not args.input and not args.batch:
        print("ERROR: Provide either --input (single file) or --batch (directory)", file=sys.stderr)
        sys.exit(1)

    if not Path(args.mapping).exists():
        print(f"ERROR: Mapping file not found: {args.mapping}", file=sys.stderr)
        print(f"  Run extract_rules.py first to generate brand-mapping.json")
        sys.exit(1)

    with open(args.mapping) as f:
        mapping = json.load(f)

    if args.batch:
        print(f"Batch converting {args.batch} → {args.output}")
        reports = batch_convert(args.batch, mapping, args.output)
        print_batch_report(reports)
    else:
        if not Path(args.input).exists():
            print(f"ERROR: Input file not found: {args.input}", file=sys.stderr)
            sys.exit(1)
        print(f"Converting: {args.input} → {args.output}")
        report = convert_deck(args.input, mapping, args.output)
        print(f"\n{report.summary()}")


if __name__ == "__main__":
    main()
